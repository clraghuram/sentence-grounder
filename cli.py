#!/usr/bin/env python3
"""
FastDocSearch v1.0
Single-file CLI tool for lightning-fast search across PDFs, DOCX, PPTX, TXT/MD and more.
Designed for quick validation of LLM responses against source documents (grounding / RAG checks).
"""

import os
import sys
import re
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Set, Tuple
import shutil

import click
from rich.console import Console
from rich.panel import Panel
from rich import print as rprint
from rich.markup import escape

# Whoosh imports
from whoosh.index import create_in, open_dir, exists_in
from whoosh.fields import Schema, TEXT, ID, STORED
from whoosh.qparser import QueryParser
from whoosh import highlight
from whoosh.analysis import StandardAnalyzer

# Document extractors
import pdfplumber
from docx import Document
from pptx import Presentation

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

console = Console()

# --------------------------- Schema ---------------------------

def get_schema() -> Schema:
    """Define Whoosh schema for document units (page/para/slide)."""
    return Schema(
        uid=ID(stored=True, unique=True),          # unique identifier e.g. /path/file.pdf:Page 5
        filepath=ID(stored=True),                  # full path for filtering/opening
        filename=STORED,                           # just basename for display
        doctype=STORED,                            # pdf / docx / pptx / txt etc.
        section=STORED,                            # "Page 12", "Paragraph 45", "Slide 3"
        content=TEXT(stored=True, analyzer=StandardAnalyzer(stoplist=None)),
        indexed_at=STORED,
    )

# --------------------------- Extractors ---------------------------

def _split_pdf_lines_into_paragraphs(lines: List[str], line_gaps: List[float], para_gap: float = 14.0) -> List[str]:
    """Split lines into paragraphs when vertical spacing indicates a break."""
    if not lines:
        return []

    paragraphs: List[str] = []
    current: List[str] = [lines[0]]

    for idx in range(1, len(lines)):
        if line_gaps[idx - 1] > para_gap:
            paragraphs.append("\n".join(current))
            current = [lines[idx]]
        else:
            current.append(lines[idx])

    paragraphs.append("\n".join(current))
    return [p.strip() for p in paragraphs if len(p.strip()) > 25]


def _extract_pdf_paragraphs(page) -> List[str]:
    """Extract paragraph-sized units from a PDF page using layout positions."""
    words = page.extract_words() or []
    if not words:
        text = (page.extract_text() or "").strip()
        if not text:
            return []
        # Fallback: blank-line splits, then per-line units for single-line PDFs.
        blocks = [b.strip() for b in re.split(r"\n\s*\n", text) if len(b.strip()) > 25]
        return blocks or ([text] if len(text) > 25 else [])

    sorted_words = sorted(words, key=lambda w: (w["top"], w["x0"]))
    lines_raw: List[List[Dict]] = []
    current_line: List[Dict] = []
    current_top: Optional[float] = None

    for word in sorted_words:
        top = word["top"]
        if current_top is None or abs(top - current_top) <= 3.0:
            current_line.append(word)
            if current_top is None:
                current_top = top
        else:
            if current_line:
                lines_raw.append(current_line)
            current_line = [word]
            current_top = top
    if current_line:
        lines_raw.append(current_line)

    line_texts: List[str] = []
    line_tops: List[float] = []
    line_bottoms: List[float] = []
    for line_words in lines_raw:
        line_texts.append(" ".join(w["text"] for w in sorted(line_words, key=lambda w: w["x0"])))
        line_tops.append(line_words[0]["top"])
        line_bottoms.append(max(w["bottom"] for w in line_words))

    gaps = [line_tops[i] - line_bottoms[i - 1] for i in range(1, len(line_texts))]
    paragraphs = _split_pdf_lines_into_paragraphs(line_texts, gaps)

    if paragraphs:
        return paragraphs

    page_text = (page.extract_text() or "").strip()
    return [page_text] if len(page_text) > 25 else []


def extract_pdf(pdf_path: Path) -> List[Dict]:
    """Extract paragraph-level units from PDF pages using pdfplumber layout."""
    results: List[Dict] = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                paragraphs = _extract_pdf_paragraphs(page)
                for para_num, text in enumerate(paragraphs, 1):
                    results.append({
                        "uid": f"{pdf_path}:{page_num}:p{para_num}",
                        "filepath": str(pdf_path),
                        "filename": pdf_path.name,
                        "doctype": "pdf",
                        "section": f"Page {page_num}, Para {para_num}",
                        "content": text,
                        "indexed_at": datetime.now().isoformat(),
                    })
    except Exception as e:
        console.print(f"[red]⚠️  PDF error in {pdf_path.name}: {e}[/red]")
    return results


def extract_docx(docx_path: Path) -> List[Dict]:
    """Extract paragraphs and table rows from DOCX."""
    results: List[Dict] = []
    try:
        doc = Document(docx_path)
        # Paragraphs
        for i, para in enumerate(doc.paragraphs, 1):
            text = para.text.strip()
            if len(text) > 25:
                results.append({
                    "uid": f"{docx_path}:p{i}",
                    "filepath": str(docx_path),
                    "filename": docx_path.name,
                    "doctype": "docx",
                    "section": f"Paragraph {i}",
                    "content": text,
                    "indexed_at": datetime.now().isoformat(),
                })
        # Tables (flattened rows)
        for t_idx, table in enumerate(doc.tables, 1):
            for r_idx, row in enumerate(table.rows, 1):
                cells = [cell.text.strip() for cell in row.cells]
                text = " | ".join(c for c in cells if c)
                if len(text) > 20:
                    results.append({
                        "uid": f"{docx_path}:t{t_idx}r{r_idx}",
                        "filepath": str(docx_path),
                        "filename": docx_path.name,
                        "doctype": "docx",
                        "section": f"Table {t_idx} Row {r_idx}",
                        "content": text,
                        "indexed_at": datetime.now().isoformat(),
                    })
    except Exception as e:
        console.print(f"[red]⚠️  DOCX error in {docx_path.name}: {e}[/red]")
    return results


def extract_pptx(pptx_path: Path) -> List[Dict]:
    """Extract text from each slide in PPTX."""
    results: List[Dict] = []
    try:
        prs = Presentation(pptx_path)
        for s_idx, slide in enumerate(prs.slides, 1):
            texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                elif hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_content = "\n".join(texts)
                results.append({
                    "uid": f"{pptx_path}:s{s_idx}",
                    "filepath": str(pptx_path),
                    "filename": pptx_path.name,
                    "doctype": "pptx",
                    "section": f"Slide {s_idx}",
                    "content": slide_content,
                    "indexed_at": datetime.now().isoformat(),
                })
    except Exception as e:
        console.print(f"[red]⚠️  PPTX error in {pptx_path.name}: {e}[/red]")
    return results


def extract_text_file(txt_path: Path) -> List[Dict]:
    """Extract paragraphs from plain text / markdown files."""
    results: List[Dict] = []
    try:
        with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        # Split on blank lines for paragraphs
        paras = [p.strip() for p in content.split("\n\n") if len(p.strip()) > 30]
        if not paras and content.strip():
            paras = [content.strip()[:8000]]  # fallback: whole file capped
        for i, para in enumerate(paras[:300]):  # safety cap per file
            results.append({
                "uid": f"{txt_path}:p{i}",
                "filepath": str(txt_path),
                "filename": txt_path.name,
                "doctype": txt_path.suffix.lower().lstrip(".") or "txt",
                "section": f"Para {i+1}",
                "content": para,
                "indexed_at": datetime.now().isoformat(),
            })
    except Exception as e:
        console.print(f"[red]⚠️  Text file error in {txt_path.name}: {e}[/red]")
    return results


def extract_xlsx(xlsx_path: Path) -> List[Dict]:
    """Extract text from Excel sheets (rows as text). Limited for practicality."""
    if not HAS_OPENPYXL:
        return []
    results: List[Dict] = []
    try:
        wb = load_workbook(xlsx_path, data_only=True, read_only=True)
        for sheet_name in wb.sheetnames[:10]:  # cap sheets
            ws = wb[sheet_name]
            row_texts: List[str] = []
            for row_idx, row in enumerate(ws.iter_rows(max_row=80), 1):  # cap rows
                cells = [str(c.value).strip() for c in row if c.value is not None]
                if cells:
                    row_texts.append(f"R{row_idx}: " + " | ".join(cells))
            if row_texts:
                content = "\n".join(row_texts)[:6000]
                results.append({
                    "uid": f"{xlsx_path}:{sheet_name}",
                    "filepath": str(xlsx_path),
                    "filename": xlsx_path.name,
                    "doctype": "xlsx",
                    "section": f"Sheet: {sheet_name}",
                    "content": content,
                    "indexed_at": datetime.now().isoformat(),
                })
        wb.close()
    except Exception as e:
        console.print(f"[red]⚠️  XLSX error in {xlsx_path.name}: {e}[/red]")
    return results


def get_extractor(file_path: Path):
    """Return the appropriate extractor function for the file extension."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf
    elif ext == ".docx":
        return extract_docx
    elif ext in (".pptx", ".ppt"):
        return extract_pptx
    elif ext in (".txt", ".md", ".rst", ".html", ".htm"):
        return extract_text_file
    elif ext in (".xlsx", ".xlsm", ".xls"):
        return extract_xlsx
    return None


def walk_and_extract(docs_dir: Path, file_types: Set[str]) -> List[Dict]:
    """Recursively find files and extract searchable units."""
    all_units: List[Dict] = []
    files: List[Path] = []
    for ext in file_types:
        files.extend(docs_dir.rglob(f"*{ext}"))
    files = sorted(set(files))  # dedupe + sort for determinism

    if not files:
        console.print("[yellow]No matching files found in the directory.[/yellow]")
        return all_units

    console.print(f"[cyan]Found {len(files)} files matching {sorted(file_types)}[/cyan]")

    for fpath in files:
        extractor = get_extractor(fpath)
        if extractor:
            console.print(f"  [blue]→[/blue] {fpath.name} ...", end=" ")
            units = extractor(fpath)
            all_units.extend(units)
            console.print(f"[green]{len(units)} units[/green]")
        else:
            console.print(f"  [dim]Skipping unsupported: {fpath.name}[/dim]")
    return all_units

# --------------------------- Indexer ---------------------------

def build_index(docs_dir: str, index_dir: str, force: bool = False, file_types_str: str = "pdf,docx,txt,md,pptx,xlsx"):
    """Build or rebuild the Whoosh index."""
    docs_path = Path(docs_dir).resolve()
    idx_path = Path(index_dir).resolve()
    idx_path.mkdir(parents=True, exist_ok=True)

    file_types: Set[str] = {f".{x.strip().lower()}" for x in file_types_str.split(",") if x.strip()}

    if exists_in(str(idx_path)) and not force:
        console.print("[yellow]Index already exists at that location. Use --force to rebuild from scratch.[/yellow]")
        return

    if force and exists_in(str(idx_path)):
        console.print("[yellow]Removing old index...[/yellow]")
        shutil.rmtree(idx_path)
        idx_path.mkdir(parents=True)

    schema = get_schema()
    ix = create_in(str(idx_path), schema)
    writer = ix.writer(limitmb=512)  # increase if you have very large docs

    console.print(f"\n[bold green]🔍 Starting indexing from:[/bold green] {docs_path}")
    units = walk_and_extract(docs_path, file_types)

    if not units:
        console.print("[red]No content extracted. Nothing to index.[/red]")
        writer.cancel()
        return

    console.print(f"\n[cyan]Writing {len(units)} text units into Whoosh index...[/cyan]")

    for unit in units:
        try:
            writer.add_document(**unit)
        except Exception as e:
            console.print(f"[red]Failed to add unit from {unit.get('filename')}: {e}[/red]")

    writer.commit()
    console.print(Panel.fit(
        f"[bold green]✅ Indexing complete![/bold green]\n\n"
        f"Index location : {idx_path}\n"
        f"Total units    : {len(units)}\n"
        f"File types     : {', '.join(sorted(file_types))}\n\n"
        f"Now you can search with:\n"
        f"  python cli.py search --query 'your phrase here' --index-dir {idx_path}",
        title="FastDocSearch",
        border_style="green"
    ))

# --------------------------- Searcher ---------------------------

def _snippet_from_hit(hit, context_chars: int = 350) -> str:
    """Build a plain-text snippet from a Whoosh hit (no Rich markup)."""
    try:
        raw_snippet = hit.highlights(
            "content",
            top=3,
            minscore=1,
            between="\n[...]\n",
        ) or (hit["content"][:context_chars] + "...")
    except Exception:
        raw_snippet = (hit.get("content") or "")[:context_chars] + "..."

    clean_text = re.sub(r"</?b[^>]*>", "", raw_snippet)
    if len(clean_text) > context_chars + 120:
        clean_text = clean_text[:context_chars] + " …"
    return clean_text


def run_query(
    searcher,
    schema,
    query_str: str,
    limit: int = 10,
    context_chars: int = 350,
) -> List[Dict]:
    """
    Run a Whoosh query and return structured hit dicts.
    Returns [] on parse errors or no matches.
    """
    parser = QueryParser("content", schema)
    try:
        q = parser.parse(query_str)
    except Exception:
        return []

    results = searcher.search(q, limit=limit)
    hits: List[Dict] = []
    for hit in results:
        content = hit.get("content") or ""
        snippet = _snippet_from_hit(hit, context_chars)
        hits.append({
            "filename": hit.get("filename", "unknown"),
            "section": hit.get("section", ""),
            "filepath": hit.get("filepath", ""),
            "doctype": (hit.get("doctype") or "").upper(),
            "score": float(hit.score),
            "snippet": snippet,
            "content": content,
        })
    return hits


def do_search(
    index_dir: str,
    query_str: str,
    limit: int = 10,
    context_chars: int = 350,
    show_score: bool = False,
    file_filter: Optional[str] = None,
):
    """Perform search and display beautiful results with snippets."""
    idx_path = Path(index_dir).resolve()
    if not exists_in(str(idx_path)):
        console.print("[red]❌ No index found at that location. Run the 'index' command first.[/red]")
        return

    ix = open_dir(str(idx_path))
    parser = QueryParser("content", ix.schema)

    try:
        parser.parse(query_str)
    except Exception as e:
        console.print(f"[red]Query parsing error: {e}[/red]")
        console.print("Tip: Use double quotes for exact phrases, e.g. \"ethics training annually\"")
        return

    with ix.searcher() as searcher:
        hits = run_query(searcher, ix.schema, query_str, limit=limit, context_chars=context_chars)

        if not hits:
            console.print("[yellow]No matches found.[/yellow]")
            console.print("Suggestions:\n"
                          "• Try an exact phrase in \"double quotes\"\n"
                          "• Use broader keywords or OR\n"
                          "• Check if the documents were indexed with the correct --file-types")
            return

        console.print(f"\n[bold green]🔎 Found {len(hits)} matching units (top {min(limit, len(hits))})[/bold green]\n")

        for i, hit in enumerate(hits, 1):
            score_str = f"  [dim](score: {hit['score']:.2f})[/dim]" if show_score else ""
            title = (
                f"[bold]{hit['filename']}[/bold]  •  {hit['section']}  •  "
                f"{hit['doctype']}{score_str}"
            )
            body = f"[dim]{escape(hit['filepath'])}[/dim]\n\n{escape(hit['snippet'])}"
            panel = Panel(
                body,
                title=title,
                border_style="blue",
                expand=False,
                padding=(0, 1),
            )
            console.print(panel)
            if i < len(hits):
                console.print()

# --------------------------- Validate: English claim extraction ---------------------------

# Common English stopwords for content-token filtering (stdlib heuristics, v1).
EN_STOPWORDS: Set[str] = {
    "a", "an", "the", "and", "or", "but", "if", "then", "else", "when", "while",
    "of", "at", "by", "for", "with", "about", "against", "between", "into", "through",
    "during", "before", "after", "above", "below", "to", "from", "up", "down", "in",
    "out", "on", "off", "over", "under", "again", "further", "once", "here", "there",
    "all", "any", "both", "each", "few", "more", "most", "other", "some", "such",
    "no", "nor", "not", "only", "own", "same", "so", "than", "too", "very", "can",
    "will", "just", "don", "should", "now", "is", "are", "was", "were", "be", "been",
    "being", "have", "has", "had", "having", "do", "does", "did", "doing", "would",
    "could", "ought", "might", "must", "shall", "this", "that", "these", "those",
    "i", "me", "my", "we", "our", "you", "your", "he", "him", "his", "she", "her",
    "it", "its", "they", "them", "their", "what", "which", "who", "whom", "whose",
    "as", "also", "however", "therefore", "thus", "hence", "moreover", "furthermore",
    "additionally", "regarding", "according", "because", "since", "although", "though",
    "whether", "where", "why", "how", "may", "per", "via", "etc", "eg", "ie",
}

# Boilerplate openings that rarely need grounding alone
_BOILERPLATE_RE = re.compile(
    r"^(note that|it is important|in conclusion|in summary|overall|"
    r"to summarize|as mentioned|as noted|please note)\b",
    re.I,
)


def _sanitize_for_whoosh(text: str) -> str:
    """Strip Whoosh query-syntax characters from free text."""
    text = re.sub(r'[:"()[\]{}~*^?\\]', " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def content_tokens(text: str) -> List[str]:
    """Lowercased alphanumeric tokens excluding English stopwords."""
    tokens = re.findall(r"[A-Za-z][A-Za-z0-9\-']{2,}", text)
    out: List[str] = []
    for t in tokens:
        low = t.lower()
        if low in EN_STOPWORDS:
            continue
        out.append(low)
    return out


def decompose_claims(text: str, max_claims: int = 20) -> List[str]:
    """
    Split one English LLM response into sentence-level claims.
    Does not parse multiple responses — entire text is one answer.
    """
    text = (text or "").strip()
    if not text:
        return []

    # Normalize whitespace but keep sentence boundaries
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)

    # Split on sentence-ending punctuation (keep abbreviations loosely intact)
    parts = re.split(r"(?<=[.!?])\s+(?=[A-Z\"'(0-9])", text)
    claims: List[str] = []
    for part in parts:
        s = part.strip().strip('"').strip()
        if len(s) < 30:
            continue
        if _BOILERPLATE_RE.match(s) and len(s) < 80:
            continue
        # Merge orphan short pieces already skipped; drop pure questions if tiny
        claims.append(s)
        if len(claims) >= max_claims:
            break
    return claims


def extract_number_phrases(text: str) -> List[str]:
    """Pull high-precision numeric anchors and short surrounding context."""
    phrases: List[str] = []
    seen: Set[str] = set()

    def _add(p: str):
        p = _sanitize_for_whoosh(p)
        if len(p) < 4:
            return
        key = p.lower()
        if key not in seen:
            seen.add(key)
            phrases.append(p)

    # Percentage ranges: 33% to 38%
    for m in re.finditer(
        r"\b\d+(?:\.\d+)?%\s+to\s+\d+(?:\.\d+)?%",
        text,
        flags=re.I,
    ):
        _add(m.group(0))

    # Rates like "0.23 patients per 100 years"
    for m in re.finditer(
        r"\b\d+(?:\.\d+)?\s+patients per 100 years(?:\s+of exposure)?",
        text,
        flags=re.I,
    ):
        _add(m.group(0))

    # "N events of ..."
    for m in re.finditer(
        r"\b\d+\s+events of [A-Za-z][A-Za-z0-9\- ]{3,40}",
        text,
        flags=re.I,
    ):
        _add(m.group(0).strip())

    # Decimal rates in parentheses e.g. (0.23 patients...)
    for m in re.finditer(
        r"\(\s*(\d+(?:\.\d+)?\s+patients per 100 years[^)]*)\)",
        text,
        flags=re.I,
    ):
        _add(m.group(1))

    # Standalone percentages only if distinctive (avoid bare "0.2%" as sole anchor)
    for m in re.finditer(r"\b\d+(?:\.\d+)?%", text):
        pct = m.group(0)
        if re.match(r"^\d{1,2}(?:\.\d+)?%$", pct) and not re.search(r"%\s+to\s+", text[m.start():m.start()+20]):
            # Prefer surrounding 4–8 words including the %
            start = max(0, m.start() - 40)
            end = min(len(text), m.end() + 40)
            ctx = text[start:end]
            words = re.findall(r"[A-Za-z0-9][A-Za-z0-9\-']*%?", ctx)
            # Find index of percent token-ish
            for i, w in enumerate(words):
                if "%" in w or w == pct.rstrip("%"):
                    lo, hi = max(0, i - 3), min(len(words), i + 4)
                    _add(" ".join(words[lo:hi]))
                    break
            else:
                if len(pct) >= 4:  # e.g. 33% not 0.2% alone — still skip short
                    pass

    return phrases


def extract_quoted_spans(text: str) -> List[str]:
    """Extract double-quoted spans suitable as exact phrases."""
    spans: List[str] = []
    for m in re.finditer(r'"([^"]{4,80})"', text):
        s = _sanitize_for_whoosh(m.group(1))
        if s:
            spans.append(s)
    return spans


def _content_phrase_windows(text: str) -> List[str]:
    """
    Build exact-phrase candidates from consecutive word windows that
    contain enough non-stopword content (stdlib heuristic).
    """
    words = re.findall(r"[A-Za-z0-9][A-Za-z0-9\-']*", text)
    if len(words) < 4:
        return []

    # score: higher content density ratio + longer rare tokens first
    scored: List[Tuple[float, float, str]] = []

    for length in range(8, 2, -1):
        for i in range(0, len(words) - length + 1):
            window = words[i : i + length]
            content = [w for w in window if w.lower() not in EN_STOPWORDS and len(w) > 2]
            # Short windows need 2 strong content words (e.g. "history of pancreatitis");
            # longer windows need 3+.
            min_content = 2 if length <= 4 else 3
            if len(content) < min_content:
                continue
            phrase = _sanitize_for_whoosh(" ".join(window))
            if len(phrase) < 10:
                continue
            ratio = len(content) / length
            rarity = sum(len(c) for c in content) / max(len(content), 1)
            # Prefer dense, rare-token windows (negated for ascending sort)
            scored.append((-ratio, -rarity, phrase))

    scored.sort()
    seen: Set[str] = set()
    out: List[str] = []
    for _, __, phrase in scored:
        key = phrase.lower()
        if key in seen:
            continue
        # Skip near-duplicates (substring of an already chosen phrase)
        if any(key in s or s in key for s in seen):
            continue
        seen.add(key)
        out.append(phrase)
        if len(out) >= 5:
            break

    # Also add consecutive original windows that are pure content bigrams/trigrams
    # preserving stopwords between (e.g. "history of pancreatitis")
    for length in (3, 4, 5):
        for i in range(0, len(words) - length + 1):
            window = words[i : i + length]
            content = [w for w in window if w.lower() not in EN_STOPWORDS and len(w) > 2]
            if len(content) < 2:
                continue
            if length == 3 and len(content) < 2:
                continue
            phrase = _sanitize_for_whoosh(" ".join(window))
            key = phrase.lower()
            if key in seen or len(phrase) < 10:
                continue
            # Prefer mid-sentence factual cores over lead-ins
            if window[0].lower() in {"according", "additionally", "regarding", "therefore"}:
                continue
            seen.add(key)
            out.append(phrase)
            if len(out) >= 6:
                return out[:6]

    return out[:5]


def plan_queries(claim: str, max_queries: int = 6) -> List[Tuple[str, str]]:
    """
    Build (strategy, query_str) candidates for one claim.
    strategy: 'exact' | 'soft'
    """
    queries: List[Tuple[str, str]] = []
    seen: Set[str] = set()

    def _add(strategy: str, q: str):
        q = q.strip()
        if not q or len(q) < 3:
            return
        key = q.lower()
        if key in seen:
            return
        seen.add(key)
        queries.append((strategy, q))

    for p in extract_number_phrases(claim):
        _add("exact", f'"{p}"')

    for p in extract_quoted_spans(claim):
        _add("exact", f'"{p}"')

    for p in _content_phrase_windows(claim):
        if len(p.split()) >= 3:
            _add("exact", f'"{p}"')

    toks = content_tokens(claim)
    # Longer tokens first; soft AND with only 2–3 terms (5-way AND is too brittle)
    ranked = sorted(set(toks), key=lambda t: (-len(t), t))
    if len(ranked) >= 2:
        _add("soft", " ".join(ranked[:3]))  # AND
    if len(ranked) >= 3:
        # Broader fallback: OR of distinctive tokens
        _add("soft", " OR ".join(ranked[:4]))
    elif len(ranked) == 1:
        _add("soft", ranked[0])

    return queries[:max_queries]


def token_overlap(claim: str, evidence: str) -> Tuple[float, int]:
    """
    Jaccard-like overlap on content tokens.
    Returns (ratio, intersection_count).
    """
    a = set(content_tokens(claim))
    b = set(content_tokens(evidence))
    if not a or not b:
        return 0.0, 0
    inter = a & b
    union = a | b
    return len(inter) / len(union), len(inter)


def score_status(strategy: str, overlap: float, inter_count: int, has_hits: bool) -> str:
    """Map retrieval quality to grounded / partial / not_found."""
    if not has_hits:
        return "not_found"
    if strategy == "exact" and inter_count >= 3:
        return "grounded"
    if overlap >= 0.40 and inter_count >= 3:
        return "grounded"
    if strategy == "exact" and inter_count >= 2:
        return "grounded"
    if inter_count >= 2 or overlap >= 0.20:
        return "partial"
    return "partial"


def search_ladder(
    searcher,
    schema,
    queries: List[Tuple[str, str]],
    limit: int = 3,
    context_chars: int = 350,
) -> Tuple[List[Dict], Optional[str], Optional[str]]:
    """
    Try queries in order; return (hits, winning_query, strategy) for first useful hit set.
    Prefers exact hits; falls through to soft.
    """
    soft_fallback: Optional[Tuple[List[Dict], str, str]] = None

    for strategy, q in queries:
        hits = run_query(searcher, schema, q, limit=limit, context_chars=context_chars)
        if not hits:
            continue
        if strategy == "exact":
            return hits, q, strategy
        if soft_fallback is None:
            soft_fallback = (hits, q, strategy)

    if soft_fallback:
        return soft_fallback
    return [], None, None


def load_response(
    response_file: Optional[str],
    text: Optional[str],
) -> Optional[str]:
    """Load a single LLM response from file, --text, or stdin."""
    if response_file:
        path = Path(response_file)
        if not path.is_file():
            console.print(f"[red]Response file not found: {path}[/red]")
            return None
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    if text is not None and text.strip():
        return text.strip()
    if not sys.stdin.isatty():
        data = sys.stdin.read().strip()
        if data:
            return data
    return None


def do_validate(
    index_dir: str,
    response_text: str,
    limit: int = 2,
    max_claims: int = 20,
    context_chars: int = 350,
    show_queries: bool = False,
):
    """Validate one English LLM response against the index (claim-level)."""
    idx_path = Path(index_dir).resolve()
    if not exists_in(str(idx_path)):
        console.print("[red]❌ No index found at that location. Run the 'index' command first.[/red]")
        return

    claims = decompose_claims(response_text, max_claims=max_claims)
    if not claims:
        console.print("[yellow]No checkable claims found in the response.[/yellow]")
        return

    ix = open_dir(str(idx_path))
    counts = {"grounded": 0, "partial": 0, "not_found": 0}

    console.print(
        f"\n[bold]Validating {len(claims)} claim(s) against[/bold] [cyan]{idx_path}[/cyan]\n"
        f"[dim]English-only · one response · stdlib phrase extraction[/dim]\n"
    )

    status_style = {
        "grounded": "green",
        "partial": "yellow",
        "not_found": "red",
    }

    with ix.searcher() as searcher:
        for i, claim in enumerate(claims, 1):
            queries = plan_queries(claim)
            hits, winning_q, strategy = search_ladder(
                searcher, ix.schema, queries, limit=limit, context_chars=context_chars
            )

            best_overlap, best_inter = 0.0, 0
            if hits:
                # Score against full stored content of best hit for robustness
                best_overlap, best_inter = token_overlap(claim, hits[0].get("content") or hits[0]["snippet"])

            status = score_status(strategy or "soft", best_overlap, best_inter, bool(hits))
            counts[status] += 1
            style = status_style[status]

            claim_display = claim if len(claim) <= 220 else claim[:217] + "..."
            header = f"[{style}][{status.upper()}][/{style}]  Claim {i}/{len(claims)}"

            lines = [f"[bold]{escape(claim_display)}[/bold]"]
            if show_queries:
                q_lines = ", ".join(f"{s}:{q}" for s, q in queries) or "(none)"
                lines.append(f"[dim]queries: {escape(q_lines)}[/dim]")
            if winning_q:
                lines.append(
                    f"[dim]matched via {strategy}: {escape(winning_q)} "
                    f"(overlap={best_overlap:.2f}, shared={best_inter})[/dim]"
                )

            if hits:
                for h in hits[:limit]:
                    lines.append(
                        f"\n[cyan]{escape(h['filename'])}[/cyan]  •  "
                        f"{escape(h['section'])}  •  {h['doctype']}\n"
                        f"[dim]{escape(h['filepath'])}[/dim]\n"
                        f"{escape(h['snippet'])}"
                    )
            else:
                lines.append("\n[dim]No supporting unit found in the index.[/dim]")

            console.print(Panel(
                "\n".join(lines),
                title=header,
                border_style=style,
                expand=False,
                padding=(0, 1),
            ))
            console.print()

    total = len(claims)
    console.print(Panel.fit(
        f"[green]grounded[/green]: {counts['grounded']}  ·  "
        f"[yellow]partial[/yellow]: {counts['partial']}  ·  "
        f"[red]not_found[/red]: {counts['not_found']}  ·  "
        f"total: {total}",
        title="Validation summary",
        border_style="cyan",
    ))


# --------------------------- CLI Commands ---------------------------

@click.group(context_settings={"help_option_names": ["-h", "--help"]})
@click.version_option("1.0.0", prog_name="FastDocSearch")
def cli():
    """
    FastDocSearch — Blazing fast local search across documents.

    Built to validate whether LLM answers are grounded in your source PDFs, DOCX, PPTX, etc.
    Index once → search in milliseconds forever.
    """
    pass


@cli.command("index")
@click.option(
    "--docs-dir", "-d", required=True,
    type=click.Path(exists=True, file_okay=False, dir_okay=True, resolve_path=True),
    help="Root directory containing your documents (scans recursively)"
)
@click.option(
    "--index-dir", "-i", default=".docsearch_index",
    type=click.Path(resolve_path=True),
    help="Where to store the Whoosh search index (default: .docsearch_index in current dir)"
)
@click.option("--force", is_flag=True, help="Delete existing index and rebuild from scratch")
@click.option(
    "--file-types", default="pdf,docx,txt,md,pptx,xlsx",
    help="Comma-separated list of extensions to include (default: pdf,docx,txt,md,pptx,xlsx)"
)
def cmd_index(docs_dir, index_dir, force, file_types):
    """Index (or re-index) all documents in a folder for super-fast searching."""
    build_index(docs_dir, index_dir, force, file_types)


@cli.command("search")
@click.option("--query", "-q", required=True, help='Search query. Use "exact phrase" for best grounding validation.')
@click.option(
    "--index-dir", "-i", default=".docsearch_index",
    type=click.Path(exists=True, resolve_path=True),
    help="Path to the previously built index directory"
)
@click.option("--limit", "-k", default=10, type=int, help="Max number of results to show (default: 10)")
@click.option("--context-chars", default=380, type=int, help="How much surrounding text to show around matches")
@click.option("--show-score", is_flag=True, help="Display BM25 relevance score for each hit")
@click.option("--file-filter", help="Optional glob to restrict results (e.g. '*policy*' or '*.pdf')")
def cmd_search(query, index_dir, limit, context_chars, show_score, file_filter):
    """Search the index and show highlighted snippets with file + page/paragraph locations."""
    do_search(index_dir, query, limit, context_chars, show_score, file_filter)


@cli.command("validate")
@click.option(
    "--response-file", "-f",
    type=click.Path(exists=True, dir_okay=False, resolve_path=True),
    help="Path to a file containing one English LLM response",
)
@click.option("--text", "-t", default=None, help="Inline LLM response text (one response)")
@click.option(
    "--index-dir", "-i", default=".docsearch_index",
    type=click.Path(resolve_path=True),
    help="Path to the previously built index directory",
)
@click.option("--limit", "-k", default=2, type=int, help="Max evidence hits per claim (default 2)")
@click.option("--max-claims", default=20, type=int, help="Max claims to process (default 20)")
@click.option("--context-chars", default=350, type=int, help="Snippet length around matches")
@click.option("--show-queries", is_flag=True, help="Show generated Whoosh queries per claim")
def cmd_validate(response_file, text, index_dir, limit, max_claims, context_chars, show_queries):
    """
    Validate one English LLM response against the index (claim-level grounding).

    The entire input is treated as a single response and split into sentences/claims.
    Uses stdlib phrase heuristics (numbers, windows, stopwords) — no spaCy/YAKE.
    """
    response_text = load_response(response_file, text)
    if not response_text:
        console.print(
            "[red]No response text provided.[/red]\n"
            "Use --response-file, --text, or pipe text on stdin."
        )
        raise SystemExit(1)
    do_validate(
        index_dir=index_dir,
        response_text=response_text,
        limit=limit,
        max_claims=max_claims,
        context_chars=context_chars,
        show_queries=show_queries,
    )


@cli.command("info")
@click.option(
    "--index-dir", "-i", default=".docsearch_index",
    type=click.Path(exists=True, resolve_path=True),
    help="Path to the index directory"
)
def cmd_info(index_dir):
    """Show statistics and info about an existing index."""
    idx_path = Path(index_dir)
    if not exists_in(str(idx_path)):
        console.print("[red]No index found at that location.[/red]")
        return

    ix = open_dir(str(idx_path))
    with ix.searcher() as searcher:
        doc_count = searcher.doc_count()

    console.print(Panel.fit(
        f"[bold]Index path:[/bold] {idx_path}\n"
        f"[bold]Searchable units:[/bold] {doc_count}\n"
        f"[bold]Created with:[/bold] FastDocSearch v1.0 (Whoosh + pdfplumber + python-docx + python-pptx)\n\n"
        f"Ready for ultra-fast searches!",
        title="📊 Index Information",
        border_style="cyan"
    ))


@cli.command("clear")
@click.option(
    "--index-dir", "-i", default=".docsearch_index",
    type=click.Path(resolve_path=True),
    help="Index directory to delete"
)
@click.confirmation_option(prompt="Are you sure you want to permanently delete this index?")
def cmd_clear(index_dir):
    """Delete an index directory (frees disk space or prepares for fresh rebuild)."""
    idx_path = Path(index_dir)
    if idx_path.exists():
        shutil.rmtree(idx_path)
        console.print(f"[green]✅ Successfully deleted index at {idx_path}[/green]")
    else:
        console.print("[yellow]No index directory found.[/yellow]")


if __name__ == "__main__":
    cli()
